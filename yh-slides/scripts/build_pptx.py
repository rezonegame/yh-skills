#!/usr/bin/env python3
# /// script
# requires-python = ">=3.10"
# dependencies = [
#     "python-pptx>=1.0.0",
#     "Pillow>=10.0.0",
#     "requests>=2.31.0",
# ]
# ///
"""
Editable PPTX export for yh-slides Path B advanced mode.

Workflow:
1. Optionally generate clean backgrounds from full-slide images.
2. Use Gemini multimodal to extract text boxes and simple style metadata.
3. Rebuild an editable PPTX with the cleaned background plus text boxes.

Example:
    uv run build_pptx.py ^
      --slides "C:\\PPTX\\demo\\images\\slide-01.png" "C:\\PPTX\\demo\\images\\slide-02.png" ^
      --output "C:\\PPTX\\demo\\output\\demo-editable.pptx" ^
      --auto-clean-bg ^
      --dump-json "C:\\PPTX\\demo\\output\\text-data.json"
"""

from __future__ import annotations

import argparse
import json
import math
import os
import re
import sys
import tempfile
from dataclasses import dataclass
from pathlib import Path
from typing import Any

import requests
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


SCRIPT_DIR = Path(__file__).resolve().parent
SKILL_DIR = SCRIPT_DIR.parent

SLIDE_WIDTH_IN = 13.333
SLIDE_HEIGHT_IN = 7.5
DEFAULT_LAYOUT = "LAYOUT_16X9"


def find_api_key(explicit_key: str | None = None) -> str:
    if explicit_key:
        return explicit_key

    env_key = os.environ.get("GEMINI_API_KEY")
    if env_key:
        return env_key

    env_candidates = [
        SKILL_DIR / ".env",
        Path.home() / ".claude" / ".env",
        Path.cwd() / ".env",
    ]
    for env_file in env_candidates:
        if not env_file.exists():
            continue
        for line in env_file.read_text(encoding="utf-8").splitlines():
            if line.startswith("GEMINI_API_KEY="):
                key = line.split("=", 1)[1].strip().strip('"').strip("'")
                if key:
                    return key

    raise SystemExit("Missing GEMINI_API_KEY. Pass --api-key or configure it in .env.")


def image_to_base64(path: Path) -> str:
    import base64

    return base64.b64encode(path.read_bytes()).decode("utf-8")


def mime_type_for(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix in {".jpg", ".jpeg"}:
        return "image/jpeg"
    if suffix == ".webp":
        return "image/webp"
    return "image/png"


def extract_json_block(text: str) -> dict[str, Any]:
    fenced = re.search(r"```(?:json)?\s*(\{.*\})\s*```", text, re.S)
    candidate = fenced.group(1) if fenced else text
    candidate = candidate.strip()
    start = candidate.find("{")
    end = candidate.rfind("}")
    if start == -1 or end == -1 or end <= start:
        raise ValueError("Model response did not contain a JSON object")
    return json.loads(candidate[start : end + 1])


def request_gemini_json(
    *,
    prompt: str,
    image_path: Path,
    api_key: str,
    model: str = "gemini-2.5-flash",
    timeout: int = 180,
    max_retries: int = 3,
) -> dict[str, Any]:
    """Send a multimodal request to Gemini and parse JSON response, with retry."""
    import time as _time
    import random as _random

    url = (
        "https://generativelanguage.googleapis.com/v1beta/"
        f"models/{model}:generateContent?key={api_key}"
    )
    payload = {
        "contents": [
            {
                "parts": [
                    {"text": prompt},
                    {
                        "inlineData": {
                            "mimeType": mime_type_for(image_path),
                            "data": image_to_base64(image_path),
                        }
                    },
                ]
            }
        ],
        "generationConfig": {
            "responseMimeType": "application/json",
        },
    }

    last_error = None
    for attempt in range(1, max_retries + 1):
        current_timeout = timeout + (attempt - 1) * 60
        try:
            response = requests.post(url, json=payload, timeout=current_timeout)

            if response.status_code == 429:
                retry_after = int(response.headers.get("Retry-After", 15))
                wait = retry_after + _random.uniform(1, 5)
                print(f"    ⭯ 429 速率限制，等待 {wait:.0f}s 后重试 ({attempt}/{max_retries})")
                _time.sleep(wait)
                continue

            if response.status_code >= 500:
                wait = (8 * (2 ** (attempt - 1))) + _random.uniform(0, 3)
                print(f"    ⭯ 服务端错误 {response.status_code}，等待 {wait:.0f}s 后重试 ({attempt}/{max_retries})")
                _time.sleep(wait)
                continue

            response.raise_for_status()
            data = response.json()
            candidates = data.get("candidates") or []
            if not candidates:
                raise ValueError("Gemini returned no candidates")
            parts = candidates[0].get("content", {}).get("parts", [])
            text = "".join(part.get("text", "") for part in parts if "text" in part)
            return extract_json_block(text)

        except requests.exceptions.Timeout as exc:
            last_error = exc
            wait = 10 + _random.uniform(0, 5)
            print(f"    ⭯ 请求超时 ({current_timeout}s)，等待 {wait:.0f}s 后重试 ({attempt}/{max_retries})")
            _time.sleep(wait)
        except requests.exceptions.ConnectionError as exc:
            last_error = exc
            wait = 10
            print(f"    ⭯ 连接错误，等待 {wait}s 后重试 ({attempt}/{max_retries})")
            _time.sleep(wait)

    raise RuntimeError(f"Gemini 请求在 {max_retries} 次重试后仍然失败: {last_error}")


def rgb_from_hex(value: str | None, fallback: str = "111111") -> RGBColor:
    color = (value or fallback).strip().lstrip("#")
    if not re.fullmatch(r"[0-9A-Fa-f]{6}", color):
        color = fallback
    return RGBColor.from_string(color.upper())


def clamp(value: float, min_value: float, max_value: float) -> float:
    return max(min_value, min(max_value, value))


def is_cjk(ch: str) -> bool:
    cp = ord(ch)
    return (
        0x4E00 <= cp <= 0x9FFF
        or 0x3400 <= cp <= 0x4DBF
        or 0x3000 <= cp <= 0x303F
        or 0xFF00 <= cp <= 0xFFEF
    )


def estimate_text_width_units(text: str, font_size: float) -> float:
    total = 0.0
    for ch in text:
        if ch == "\n":
            continue
        if ch.isspace():
            total += font_size * 0.33
        elif is_cjk(ch):
            total += font_size * 1.0
        elif ch.isupper():
            total += font_size * 0.62
        else:
            total += font_size * 0.52
    return total


def estimate_wrapped_lines(text: str, font_size: float, box_width_pt: float) -> int:
    lines = 1
    for paragraph in text.splitlines() or [""]:
        if not paragraph:
            lines += 1
            continue
        width = estimate_text_width_units(paragraph, font_size)
        lines += max(1, math.ceil(width / max(box_width_pt, 1.0))) - 1
    return max(lines, 1)


def estimate_fitting_font_size(text: str, box_width_pt: float, box_height_pt: float, upper_bound: float) -> float:
    low = 8.0
    high = max(low, upper_bound)
    best = low
    while high - low > 0.5:
        mid = (low + high) / 2
        lines = estimate_wrapped_lines(text, mid, box_width_pt)
        estimated_height = lines * mid * 1.22
        if estimated_height <= box_height_pt:
            best = mid
            low = mid
        else:
            high = mid
    return round(best, 1)


def normalize_segments(item: dict[str, Any]) -> list[dict[str, Any]]:
    segments = item.get("colored_segments")
    if isinstance(segments, list) and segments:
        return [
            {
                "text": str(segment.get("text", "")),
                "color": segment.get("color") or item.get("color") or "#111111",
                "bold": bool(segment.get("bold", item.get("bold", False))),
                "italic": bool(segment.get("italic", item.get("italic", False))),
            }
            for segment in segments
            if str(segment.get("text", "")).strip()
        ]
    text = str(item.get("text", ""))
    return [
        {
            "text": text,
            "color": item.get("color") or "#111111",
            "bold": bool(item.get("bold", False)),
            "italic": bool(item.get("italic", False)),
        }
    ]


@dataclass
class TextItem:
    text: str
    x: float
    y: float
    w: float
    h: float
    align: str = "left"
    color: str = "#111111"
    bold: bool = False
    italic: bool = False
    font_size_pt: float | None = None
    font_family: str = "sans"
    colored_segments: list[dict[str, Any]] | None = None


def parse_text_items(slide_data: dict[str, Any]) -> list[TextItem]:
    items: list[TextItem] = []
    for raw in slide_data.get("items", []):
        bbox = raw.get("bbox") or {}
        text = str(raw.get("text", "")).strip()
        if not text:
            continue
        items.append(
            TextItem(
                text=text,
                x=float(bbox.get("x", 0.0)),
                y=float(bbox.get("y", 0.0)),
                w=float(bbox.get("w", 0.0)),
                h=float(bbox.get("h", 0.0)),
                align=str(raw.get("align", "left")).lower(),
                color=str(raw.get("color", "#111111")),
                bold=bool(raw.get("bold", False)),
                italic=bool(raw.get("italic", False)),
                font_size_pt=float(raw["font_size_pt"]) if raw.get("font_size_pt") else None,
                font_family=str(raw.get("font_family", "sans")),
                colored_segments=normalize_segments(raw),
            )
        )
    return items


def extraction_prompt() -> str:
    return """
你是一名PPT版面分析助手。请识别这张幻灯片图片中的所有“可编辑文字块”，输出严格 JSON。

要求：
1. 只提取需要在PPT里重建的文字，不提取纯装饰纹理。
2. 按视觉阅读顺序输出 items。
3. bbox 使用相对于整张图的归一化坐标，范围 0 到 1。
4. 每个 item 输出：
   - text: 完整文字
   - bbox: {x, y, w, h}
   - align: left / center / right
   - color: #RRGGBB
   - bold: true/false
   - italic: true/false
   - font_size_pt: 估计字号（pt）
   - font_family: sans / serif / mono
   - colored_segments: 可选；如果同一文字块存在多色文字，按顺序拆分片段
5. 不要输出解释，不要输出 markdown，只输出 JSON。

JSON schema:
{
  "canvas": {"width": 2048, "height": 1152},
  "items": [
    {
      "text": "标题",
      "bbox": {"x": 0.1, "y": 0.12, "w": 0.5, "h": 0.1},
      "align": "left",
      "color": "#111111",
      "bold": true,
      "italic": false,
      "font_size_pt": 36,
      "font_family": "sans",
      "colored_segments": [
        {"text": "增长", "color": "#111111", "bold": true},
        {"text": "23%", "color": "#BF9A4A", "bold": true}
      ]
    }
  ]
}
""".strip()


def map_font_family(value: str) -> str:
    family = value.lower()
    if family == "serif":
        return "Times New Roman"
    if family == "mono":
        return "Consolas"
    return "Microsoft YaHei"


def map_align(value: str) -> PP_ALIGN:
    if value == "center":
        return PP_ALIGN.CENTER
    if value == "right":
        return PP_ALIGN.RIGHT
    return PP_ALIGN.LEFT


def add_background(slide, image_path: Path) -> None:
    slide.shapes.add_picture(str(image_path), 0, 0, Inches(SLIDE_WIDTH_IN), Inches(SLIDE_HEIGHT_IN))


def add_text_item(slide, item: TextItem) -> None:
    left = Inches(SLIDE_WIDTH_IN * clamp(item.x, 0.0, 1.0))
    top = Inches(SLIDE_HEIGHT_IN * clamp(item.y, 0.0, 1.0))
    width = Inches(SLIDE_WIDTH_IN * clamp(item.w, 0.001, 1.0))
    height = Inches(SLIDE_HEIGHT_IN * clamp(item.h, 0.001, 1.0))

    textbox = slide.shapes.add_textbox(left, top, width, height)
    textbox.line.fill.background()
    textbox.fill.background()

    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    text_frame.margin_left = 0
    text_frame.margin_right = 0
    text_frame.margin_top = 0
    text_frame.margin_bottom = 0
    text_frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = map_align(item.align)

    box_width_pt = width.pt
    box_height_pt = height.pt
    suggested_upper = item.font_size_pt or min(box_height_pt * 0.8, 44.0)
    fitted_size = estimate_fitting_font_size(item.text, box_width_pt, box_height_pt, suggested_upper)
    font_name = map_font_family(item.font_family)

    first = True
    for segment in item.colored_segments or []:
        run = paragraph.runs[0] if first and paragraph.runs else paragraph.add_run()
        if first and not paragraph.runs:
            run = paragraph.add_run()
        run.text = segment["text"]
        run.font.name = font_name
        run.font.size = Pt(fitted_size)
        run.font.bold = bool(segment.get("bold", item.bold))
        run.font.italic = bool(segment.get("italic", item.italic))
        run.font.color.rgb = rgb_from_hex(segment.get("color"), item.color.lstrip("#"))
        first = False

    if first:
        paragraph.text = item.text
        paragraph.font.name = font_name
        paragraph.font.size = Pt(fitted_size)
        paragraph.font.bold = item.bold
        paragraph.font.italic = item.italic
        paragraph.font.color.rgb = rgb_from_hex(item.color)


def ensure_clean_backgrounds(
    slide_paths: list[Path],
    bg_dir: Path,
    api_key: str,
    force: bool = False,
) -> list[Path]:
    """为每页生成干净背景。单页失败时降级使用原图，不中断整个流程。"""
    sys.path.insert(0, str(SCRIPT_DIR))
    import generate_image as gi  # noqa: PLC0415

    bg_dir.mkdir(parents=True, exist_ok=True)
    results: list[Path] = []
    failed_pages: list[str] = []

    for index, slide_path in enumerate(slide_paths, start=1):
        bg_path = bg_dir / f"{slide_path.stem}-bg.png"

        if not force and bg_path.exists():
            print(f"  [{index}/{len(slide_paths)}] 已存在: {bg_path.name}")
            results.append(bg_path)
            continue

        print(f"  [{index}/{len(slide_paths)}] 清理背景: {slide_path.name}")
        try:
            ok = gi.clean_background(
                image_path=str(slide_path),
                output_path=str(bg_path),
                api_key=api_key,
                max_retries=3,
            )
        except Exception as exc:
            ok = False
            print(f"  ⚠️  clean-bg 异常: {exc}")

        if ok and bg_path.exists():
            results.append(bg_path)
        else:
            # 降级：使用原图作为背景（文字会被重建的文本框叠加覆盖）
            print(f"  ⚠️  第 {index} 页背景清理失败，降级使用原图作为背景")
            print(f"      提示: 该页可能出现文字重叠，建议后续手动补清或重试")
            results.append(slide_path)
            failed_pages.append(slide_path.name)

    if failed_pages:
        print(f"\n⚠️  共 {len(failed_pages)} 页背景清理失败，已降级使用原图:")
        for name in failed_pages:
            print(f"    - {name}")
        print("  建议: 稍后用 generate_image.py clean-bg 单独重试这些页，")
        print("  然后用 --backgrounds 参数指定已清理的背景重新构建。\n")

    return results


def build_presentation(
    *,
    slide_paths: list[Path],
    background_paths: list[Path],
    text_data: list[dict[str, Any]],
    output_path: Path,
) -> None:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_IN)
    prs.slide_height = Inches(SLIDE_HEIGHT_IN)

    for slide_image, background_image, slide_data in zip(slide_paths, background_paths, text_data):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_background(slide, background_image)
        for item in parse_text_items(slide_data):
            add_text_item(slide, item)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


def load_text_data(path: Path) -> list[dict[str, Any]]:
    data = json.loads(path.read_text(encoding="utf-8"))
    if isinstance(data, dict) and "slides" in data:
        return list(data["slides"])
    if isinstance(data, list):
        return data
    raise ValueError("Text data JSON must be a list or an object with a 'slides' field")


def dump_text_data(path: Path, slides_data: list[dict[str, Any]]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    payload = {"slides": slides_data}
    path.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")


def collect_slide_data(slide_paths: list[Path], api_key: str, model: str) -> list[dict[str, Any]]:
    """Collect text extraction data for each slide. Single page failure returns empty items."""
    prompt = extraction_prompt()
    results: list[dict[str, Any]] = []
    failed_pages: list[str] = []

    for index, slide_path in enumerate(slide_paths, start=1):
        print(f"[{index}/{len(slide_paths)}] 提取文字框: {slide_path.name}")
        try:
            data = request_gemini_json(
                prompt=prompt,
                image_path=slide_path,
                api_key=api_key,
                model=model,
            )
            results.append(data)
        except Exception as exc:
            print(f"  ⚠️  第 {index} 页文字提取失败: {exc}")
            print(f"      该页将作为纯背景图输出（无可编辑文字框）")
            results.append({"items": []})
            failed_pages.append(slide_path.name)

    if failed_pages:
        print(f"\n⚠️  共 {len(failed_pages)} 页文字提取失败（将作为纯背景图）:")
        for name in failed_pages:
            print(f"    - {name}")

    return results


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Build editable PPTX from Path B full-slide images.")
    parser.add_argument("--slides", nargs="+", required=True, help="Full-slide image paths in order")
    parser.add_argument("--output", "-o", required=True, help="Output PPTX path")
    parser.add_argument("--api-key", help="Gemini API key")
    parser.add_argument("--model", default="gemini-2.5-flash", help="Gemini model for text extraction")
    parser.add_argument("--backgrounds", nargs="*", help="Optional cleaned background image paths in order")
    parser.add_argument("--bg-dir", help="Directory for auto-generated clean backgrounds")
    parser.add_argument("--auto-clean-bg", action="store_true", help="Generate cleaned backgrounds automatically")
    parser.add_argument("--force-clean-bg", action="store_true", help="Regenerate clean backgrounds even if they exist")
    parser.add_argument("--text-data", help="Existing extracted text-data JSON")
    parser.add_argument("--dump-json", help="Write extracted text-data JSON for review/reuse")
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    slide_paths = [Path(p).resolve() for p in args.slides]
    for path in slide_paths:
        if not path.exists():
            raise SystemExit(f"Slide image not found: {path}")

    output_path = Path(args.output).resolve()
    api_key = find_api_key(args.api_key)

    if args.text_data:
        text_data = load_text_data(Path(args.text_data))
    else:
        text_data = collect_slide_data(slide_paths, api_key, args.model)
        if args.dump_json:
            dump_text_data(Path(args.dump_json).resolve(), text_data)

    if len(text_data) != len(slide_paths):
        raise SystemExit("The number of extracted slide records does not match the number of slide images.")

    if args.backgrounds:
        background_paths = [Path(p).resolve() for p in args.backgrounds]
    elif args.auto_clean_bg:
        if args.bg_dir:
            bg_dir = Path(args.bg_dir).resolve()
        else:
            bg_dir = output_path.parent / "backgrounds"
        background_paths = ensure_clean_backgrounds(slide_paths, bg_dir, api_key, force=args.force_clean_bg)
    else:
        background_paths = slide_paths

    if len(background_paths) != len(slide_paths):
        raise SystemExit("The number of background images must match the number of slide images.")

    build_presentation(
        slide_paths=slide_paths,
        background_paths=background_paths,
        text_data=text_data,
        output_path=output_path,
    )
    print(f"Editable PPTX saved: {output_path}")


if __name__ == "__main__":
    main()
