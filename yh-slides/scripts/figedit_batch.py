#!/usr/bin/env python3
"""Batch bridge from yh-slides 2B-R to the independent FigEdit skill."""

from __future__ import annotations

import argparse
import json
import os
import shutil
import subprocess
import sys
import tempfile
import zipfile
from pathlib import Path
from typing import Any


DEFAULT_FIGEDIT_ROOT = Path(r"C:\Users\wudao\OneDrive\skills\figedit")
REQUIRED_MODULES = {
    "cv2": "opencv-python>=4.9",
    "paddleocr": "paddleocr>=3.7",
    "paddle": "paddlepaddle>=3.3",
    "numpy": "numpy>=1.24",
    "PIL": "Pillow>=10.0",
    "scipy": "scipy>=1.10",
    "matplotlib": "matplotlib>=3.7",
    "latex2mathml": "latex2mathml>=3.81",
    "lxml": "lxml>=5.0",
    "pptx": "python-pptx>=1.0",
}
REQUIRED_REVIEW_CHECKS = (
    "major_structure",
    "text_complete",
    "connectors_correct",
    "distinctive_assets_preserved",
)


def _figedit_root(explicit: str | None = None) -> Path:
    root = Path(explicit or os.environ.get("FIGEDIT_SKILL_DIR") or DEFAULT_FIGEDIT_ROOT).resolve()
    return root


def _scripts(root: Path) -> Path:
    return root / "scripts"


def _read_json(path: Path) -> dict[str, Any]:
    return json.loads(path.read_text(encoding="utf-8"))


def _write_json(path: Path, payload: Any) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(payload, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def _run(cmd: list[str], *, cwd: Path | None = None) -> None:
    print(">", subprocess.list2cmdline(cmd))
    subprocess.run(cmd, cwd=str(cwd) if cwd else None, check=True)


def preflight(root: Path, *, check_imports: bool = True) -> dict[str, Any]:
    required_files = [
        root / "SKILL.md",
        _scripts(root) / "prepare_measurements.py",
        _scripts(root) / "compose_svg_package.py",
        _scripts(root) / "svg_to_pptx" / "pptx_builder.py",
        _scripts(root) / "pptx_math.py",
    ]
    missing_files = [str(path) for path in required_files if not path.exists()]
    missing_modules: list[dict[str, str]] = []
    if check_imports:
        for module, package in REQUIRED_MODULES.items():
            proc = subprocess.run(
                [sys.executable, "-c", f"import {module}"],
                stdout=subprocess.DEVNULL,
                stderr=subprocess.DEVNULL,
            )
            if proc.returncode:
                missing_modules.append({"module": module, "package": package})
    result = {
        "status": "ok" if not missing_files and not missing_modules else "failed",
        "figedit_root": str(root),
        "python": sys.version.split()[0],
        "missing_files": missing_files,
        "missing_modules": missing_modules,
        "install_command": f'"{sys.executable}" -m pip install -r "{root / "requirements.txt"}"',
    }
    return result


def _load_batch(project_dir: Path) -> dict[str, Any]:
    path = project_dir / "reconstruction" / "batch.json"
    if not path.exists():
        raise SystemExit(f"Batch manifest not found: {path}")
    return _read_json(path)


def init_batch(project_dir: Path, slide_paths: list[Path], root: Path) -> None:
    reconstruction = project_dir / "reconstruction"
    reconstruction.mkdir(parents=True, exist_ok=True)
    pages = []
    for index, source in enumerate(slide_paths, start=1):
        source = source.resolve()
        if not source.exists():
            raise SystemExit(f"Input image not found: {source}")
        page_id = f"page-{index:02d}"
        page_dir = reconstruction / page_id
        page_dir.mkdir(parents=True, exist_ok=True)
        local_source = page_dir / f"source{source.suffix.lower() or '.png'}"
        shutil.copy2(source, local_source)
        pages.append(
            {
                "id": page_id,
                "order": index,
                "original_source": str(source),
                "source": str(local_source),
                "measurements": str(page_dir / "measurements"),
                "manifest": str(page_dir / "manifest.json"),
                "out": str(page_dir / "out"),
            }
        )
    _write_json(
        reconstruction / "batch.json",
        {
            "version": 1,
            "engine": "figedit",
            "figedit_root": str(root),
            "project_dir": str(project_dir.resolve()),
            "pages": pages,
        },
    )
    print(f"Initialized {len(pages)} page(s): {reconstruction / 'batch.json'}")


def measure_batch(
    project_dir: Path,
    root: Path,
    *,
    lang: str | None,
    gpu: bool,
    ocr_profile: str,
) -> None:
    check = preflight(root)
    if check["status"] != "ok":
        print(json.dumps(check, ensure_ascii=False, indent=2))
        raise SystemExit(2)
    batch = _load_batch(project_dir)
    for page in batch["pages"]:
        cmd = [
            sys.executable,
            str(_scripts(root) / "prepare_measurements.py"),
            page["source"],
            "--out",
            page["measurements"],
            "--ocr-profile",
            ocr_profile,
        ]
        if lang:
            cmd.extend(["--lang", lang])
        if gpu:
            cmd.append("--gpu")
        _run(cmd, cwd=project_dir)
        draft = Path(page["measurements"]) / "draft_manifest.json"
        if draft.exists():
            shutil.copy2(draft, Path(page["manifest"]).with_name("manifest.draft.json"))
    print("Measurement evidence is ready. Author and review each page manifest.json before compose.")


def compose_batch(project_dir: Path, root: Path) -> None:
    check = preflight(root)
    if check["status"] != "ok":
        print(json.dumps(check, ensure_ascii=False, indent=2))
        raise SystemExit(2)
    batch = _load_batch(project_dir)
    missing = [page["id"] for page in batch["pages"] if not Path(page["manifest"]).exists()]
    if missing:
        raise SystemExit(f"Missing model-authored manifest.json for: {', '.join(missing)}")
    failed: list[str] = []
    for page in batch["pages"]:
        out_dir = Path(page["out"])
        try:
            _run(
                [
                    sys.executable,
                    str(_scripts(root) / "compose_svg_package.py"),
                    page["manifest"],
                    "--out",
                    page["out"],
                ],
                cwd=project_dir,
            )
            (out_dir / "compose_error.json").unlink(missing_ok=True)
        except subprocess.CalledProcessError as exc:
            failed.append(page["id"])
            _write_json(
                out_dir / "compose_error.json",
                {"status": "failed", "returncode": exc.returncode, "command": exc.cmd},
            )
    result = write_status(project_dir)
    if failed or result["status"] != "ok":
        raise SystemExit(f"FigEdit compose/quality gate failed for: {', '.join(failed) or 'review required'}")


def _page_gate(page: dict[str, Any]) -> dict[str, Any]:
    out_dir = Path(page["out"])
    manifest_path = out_dir / "manifest.json"
    required = [
        out_dir / "editable.svg",
        out_dir / "editable.pptx",
        out_dir / "quality_report.md",
        out_dir / "editability_report.md",
        manifest_path,
    ]
    reasons = [f"missing:{path.name}" for path in required if not path.exists()]
    if (out_dir / "compose_error.json").exists():
        reasons.append("compose:failed")
    manifest: dict[str, Any] = {}
    if manifest_path.exists():
        manifest = _read_json(manifest_path)
        gates = manifest.get("quality_gates") or {}
        for gate_name in ("xml_editable", "xml_embedded", "pptx_export", "formula_text_leakage"):
            status = (gates.get(gate_name) or {}).get("status")
            if status not in {"ok", "skipped"}:
                reasons.append(f"{gate_name}:{status or 'missing'}")
        math_gate = gates.get("pptx_math_export") or {}
        if int(math_gate.get("attempted_count") or 0) and math_gate.get("status") != "ok":
            reasons.append(f"pptx_math_export:{math_gate.get('status') or 'missing'}")
        editability = gates.get("editability") or {}
        if editability.get("status") != "ok":
            reasons.append(f"editability:{editability.get('status') or 'missing'}")
        review = manifest.get("delivery_review") or {}
        if review.get("status") != "approved":
            reasons.append("delivery_review:not-approved")
        for check_name in REQUIRED_REVIEW_CHECKS:
            if review.get(check_name) is not True:
                reasons.append(f"delivery_review:{check_name}")
    return {
        "id": page["id"],
        "status": "ok" if not reasons else "failed",
        "reasons": reasons,
        "out": str(out_dir),
    }


def write_status(project_dir: Path) -> dict[str, Any]:
    batch = _load_batch(project_dir)
    pages = [_page_gate(page) for page in batch["pages"]]
    failed = [page for page in pages if page["status"] != "ok"]
    reconstruction = project_dir / "reconstruction"
    _write_json(reconstruction / "failed_pages.json", failed)
    lines = [
        "# 2B-R FigEdit Reconstruction Summary",
        "",
        f"- Overall status: `{'failed' if failed else 'ok'}`",
        f"- Pages: {len(pages)}",
        f"- Passed: {len(pages) - len(failed)}",
        f"- Failed: {len(failed)}",
        "",
        "## Page Results",
        "",
    ]
    for page in pages:
        reasons = ", ".join(page["reasons"]) if page["reasons"] else "all gates passed"
        lines.append(f"- {page['id']}: `{page['status']}` — {reasons}")
    lines.extend(
        [
            "",
            "## Delivery Rule",
            "",
            "- The combined editable PPTX is generated only when every page passes automated gates and explicit semantic review.",
            "- A failed run keeps measurements, manifests, previews, assets, and reports for repair.",
        ]
    )
    (reconstruction / "summary_report.md").write_text("\n".join(lines) + "\n", encoding="utf-8")
    result = {"status": "failed" if failed else "ok", "pages": pages}
    print(json.dumps(result, ensure_ascii=False, indent=2))
    return result


def _patch_math_by_slide(pptx_path: Path, manifest_paths: list[Path], scripts_dir: Path) -> list[dict[str, Any]]:
    sys.path.insert(0, str(scripts_dir))
    import pptx_math  # type: ignore

    prepared = [pptx_math.prepare_editable_math(path) for path in manifest_paths]
    reports: list[dict[str, Any]] = []
    if not any(item.get("prepared") for item in prepared):
        return [
            {
                "status": item.get("status", "skipped"),
                "attempted_count": item.get("attempted_count", 0),
                "editable_count": 0,
                "failures": item.get("failures", []),
            }
            for item in prepared
        ]

    from lxml import etree

    with tempfile.TemporaryDirectory(prefix="yh_figedit_math_") as tmp:
        extract_dir = Path(tmp) / "pptx"
        with zipfile.ZipFile(pptx_path, "r") as archive:
            archive.extractall(extract_dir)
        parser = etree.XMLParser(resolve_entities=False, no_network=True, remove_blank_text=False)
        for index, item in enumerate(prepared, start=1):
            slide_path = extract_dir / "ppt" / "slides" / f"slide{index}.xml"
            tree = etree.parse(str(slide_path), parser=parser)
            root = tree.getroot()
            sp_tree = root.find(".//p:spTree", namespaces=pptx_math.NS)
            if sp_tree is None:
                raise RuntimeError(f"{slide_path.name} has no p:spTree")
            next_id = pptx_math._max_shape_id(root) + 1
            inserted = 0
            for formula in item.get("prepared", []):
                shape_xml = pptx_math._equation_shape_xml(next_id, formula["element"], formula["omml_para_xml"])
                sp_tree.append(etree.fromstring(shape_xml.encode("utf-8"), parser=parser))
                next_id += 1
                inserted += 1
            tree.write(str(slide_path), encoding="UTF-8", xml_declaration=True, standalone=True)
            failures = item.get("failures", [])
            status = "ok" if inserted == item.get("attempted_count", 0) and not failures else "review"
            if item.get("attempted_count", 0) and inserted == 0:
                status = "failed"
            reports.append(
                {
                    "status": status if item.get("attempted_count", 0) else "skipped",
                    "attempted_count": item.get("attempted_count", 0),
                    "editable_count": inserted,
                    "failures": failures,
                }
            )
        patched = Path(tmp) / "patched.pptx"
        with zipfile.ZipFile(patched, "w", zipfile.ZIP_DEFLATED) as archive:
            for file_path in extract_dir.rglob("*"):
                if file_path.is_file():
                    archive.write(file_path, file_path.relative_to(extract_dir))
        shutil.copy2(patched, pptx_path)
    return reports


def _validate_combined_pptx(pptx_path: Path, manifest_paths: list[Path]) -> dict[str, Any]:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    presentation = Presentation(str(pptx_path))
    failures: list[str] = []
    slides: list[dict[str, Any]] = []
    if len(presentation.slides) != len(manifest_paths):
        failures.append(f"slide_count:{len(presentation.slides)}!={len(manifest_paths)}")

    def walk(shapes):
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from walk(shape.shapes)
            else:
                yield shape

    for index, (slide, manifest_path) in enumerate(zip(presentation.slides, manifest_paths), start=1):
        manifest = _read_json(manifest_path)
        leaves = list(walk(slide.shapes))
        actual_texts = [
            shape.text
            for shape in leaves
            if getattr(shape, "has_text_frame", False) and shape.text.strip()
        ]
        actual_pictures = sum(1 for shape in leaves if shape.shape_type == MSO_SHAPE_TYPE.PICTURE)
        expected_texts = sum(
            1 for element in manifest.get("elements", []) if element.get("type") == "text"
        )
        expected_pictures = sum(
            1 for element in manifest.get("elements", []) if element.get("type") == "image"
        )
        page_failures: list[str] = []
        if len(actual_texts) < expected_texts:
            page_failures.append(f"editable_texts:{len(actual_texts)}<{expected_texts}")
        if actual_pictures < expected_pictures:
            page_failures.append(f"pictures:{actual_pictures}<{expected_pictures}")
        failures.extend(f"slide-{index}:{failure}" for failure in page_failures)
        slides.append(
            {
                "slide": index,
                "leaf_shapes": len(leaves),
                "editable_text_count": len(actual_texts),
                "expected_text_count": expected_texts,
                "picture_count": actual_pictures,
                "expected_picture_count": expected_pictures,
                "status": "ok" if not page_failures else "failed",
                "failures": page_failures,
            }
        )
    return {"status": "ok" if not failures else "failed", "failures": failures, "slides": slides}


def assemble_batch(project_dir: Path, root: Path, output: Path) -> None:
    check = preflight(root)
    if check["status"] != "ok":
        print(json.dumps(check, ensure_ascii=False, indent=2))
        raise SystemExit(2)
    status = write_status(project_dir)
    if status["status"] != "ok":
        raise SystemExit("Quality gate failed. Combined PPTX was not generated.")
    batch = _load_batch(project_dir)
    svg_files = [Path(page["out"]) / "editable.svg" for page in batch["pages"]]
    manifests = [Path(page["out"]) / "manifest.json" for page in batch["pages"]]
    scripts_dir = _scripts(root)
    sys.path.insert(0, str(scripts_dir))
    from svg_to_pptx.pptx_builder import create_pptx_with_native_svg  # type: ignore

    output = output.resolve()
    output.parent.mkdir(parents=True, exist_ok=True)
    temp_output = output.with_name(f".{output.stem}.building{output.suffix}")
    if temp_output.exists():
        temp_output.unlink()
    ok = create_pptx_with_native_svg(
        svg_files=svg_files,
        output_path=temp_output,
        use_native_shapes=True,
        use_compat_mode=False,
        conversion_trace_path=output.with_suffix(output.suffix + ".trace.json"),
    )
    if not ok or not temp_output.exists():
        raise SystemExit("FigEdit native multi-page PPTX export failed.")
    math_reports = _patch_math_by_slide(temp_output, manifests, scripts_dir)
    failed_math = [index + 1 for index, item in enumerate(math_reports) if item["status"] not in {"ok", "skipped"}]
    _write_json(output.with_suffix(output.suffix + ".math_report.json"), {"slides": math_reports})
    if failed_math:
        temp_output.unlink(missing_ok=True)
        raise SystemExit(f"Editable formula gate failed on slide(s): {failed_math}")
    validation = _validate_combined_pptx(temp_output, manifests)
    _write_json(output.with_suffix(output.suffix + ".validation.json"), validation)
    if validation["status"] != "ok":
        temp_output.unlink(missing_ok=True)
        raise SystemExit(f"Combined PPTX object validation failed: {validation['failures']}")
    os.replace(temp_output, output)
    print(f"Editable reconstruction saved: {output}")


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="yh-slides 2B-R bridge to independent FigEdit")
    parser.add_argument("--figedit-root", help="Override FigEdit skill directory")
    sub = parser.add_subparsers(dest="command", required=True)

    sub.add_parser("preflight")

    init = sub.add_parser("init")
    init.add_argument("--project-dir", required=True, type=Path)
    init.add_argument("--slides", required=True, nargs="+", type=Path)

    measure = sub.add_parser("measure")
    measure.add_argument("--project-dir", required=True, type=Path)
    measure.add_argument("--lang")
    measure.add_argument("--gpu", action="store_true")
    measure.add_argument(
        "--ocr-profile",
        default="auto",
        choices=["auto", "v6_medium", "v6_small", "v6_tiny", "v5_mobile"],
    )

    compose = sub.add_parser("compose")
    compose.add_argument("--project-dir", required=True, type=Path)

    status = sub.add_parser("status")
    status.add_argument("--project-dir", required=True, type=Path)

    assemble = sub.add_parser("assemble")
    assemble.add_argument("--project-dir", required=True, type=Path)
    assemble.add_argument("--output", required=True, type=Path)
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    root = _figedit_root(args.figedit_root)
    if args.command == "preflight":
        result = preflight(root)
        print(json.dumps(result, ensure_ascii=False, indent=2))
        raise SystemExit(0 if result["status"] == "ok" else 2)
    project_dir = args.project_dir.resolve()
    if args.command == "init":
        init_batch(project_dir, args.slides, root)
    elif args.command == "measure":
        measure_batch(project_dir, root, lang=args.lang, gpu=args.gpu, ocr_profile=args.ocr_profile)
    elif args.command == "compose":
        compose_batch(project_dir, root)
    elif args.command == "status":
        result = write_status(project_dir)
        raise SystemExit(0 if result["status"] == "ok" else 3)
    elif args.command == "assemble":
        assemble_batch(project_dir, root, args.output)


if __name__ == "__main__":
    main()
